from typing import List

from langchain_core.tools import tool
from pptx import Presentation
from pydantic import BaseModel, Field


class SlideContent(BaseModel):
    header: str = Field(description="該頁簡報的標題")
    items: List[str] = Field(description="該頁簡報的條列重點")


class PPTInput(BaseModel):
    filename: str = Field(description="輸出的檔案名稱，不需要包含副檔名")
    slides: List[SlideContent] = Field(description="要生成的簡報頁面列表")


@tool("generate_ppt", args_schema=PPTInput)
def generate_ppt(filename: str, slides: List[SlideContent]) -> str:
    """用來生成策略總結 PPT 的工具。
    能夠將內容分為多頁投影片，每頁包含標題與重點列表。
    """
    try:
        prs = Presentation()

        for i, slide_content in enumerate(slides):
            # 第一頁使用 Title Slide (layout index 0)，其他使用 Title and Content (layout index 1)
            if i == 0:
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)

                # 設定標題 (Title)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_content.header

                # 設定副標題 (Subtitle) - 位於 index 1
                if len(slide.shapes.placeholders) > 1:
                    subtitle = slide.shapes.placeholders[1].text_frame
                    subtitle.clear()
                    for item in slide_content.items:
                        p = subtitle.add_paragraph()
                        p.text = item
            else:
                # 其他頁面使用 Title and Content
                layout_index = 1
                slide_layout = prs.slide_layouts[layout_index]
                slide = prs.slides.add_slide(slide_layout)

                # 設定標題
                if slide.shapes.title:
                    slide.shapes.title.text = slide_content.header

                # 設定內容
                # 檢查是否有副標題/內容框 placeholder
                if len(slide.shapes.placeholders) > 1:
                    tf = slide.shapes.placeholders[1].text_frame
                    # 清除預設內容(如果有的話)
                    tf.clear()

                    for item in slide_content.items:
                        p = tf.add_paragraph()
                        p.text = item
                        p.level = 0  # 設定縮排層級

        # 儲存
        full_path = f"{filename}.pptx"
        prs.save(full_path)

        return f"成功生成多頁簡報檔案：{full_path}"
    except Exception as e:
        return f"生成失敗：{str(e)}"
